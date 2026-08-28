"""P2-D4 验证：PPTX 导出生成器（走真实 HTTP 导出路由 → 下载 → python-pptx 打开校验）

断言：
  1. POST /api/export/report (format=pptx) → task.status=completed + download_url
  2. GET download_url → 200 且扩展名为 .pptx
  3. python-pptx 打开：幻灯片数 >= 2（封面 + 至少 1 章节）
"""
import sys, os, json, asyncio, io, time

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.environ.setdefault("OMP_NUM_THREADS", "1")

BASE = "http://localhost:8000"

SAMPLE_MD = """## 方案背景

制造业客户面临设备停机率高、人工巡检成本大的问题。

## 方案架构

- 工业物联网平台统一接入设备数据
- AI 预测性维护模型提前识别故障风险
- 联动工单系统自动派发维修任务

## 实施路径

第一阶段：数据接入与模型训练；第二阶段：试点运行与调优；第三阶段：规模化推广。
"""


async def main():
    print("=== P2-D4 PPTX 导出验证（HTTP 链路） ===")
    import httpx

    async with httpx.AsyncClient(timeout=180) as c:
        # 1. 发起 PPTX 导出
        t0 = time.time()
        r = await c.post(f"{BASE}/api/export/report", json={
            "report_type": "solution",
            "format": "pptx",
            "content": SAMPLE_MD,
            "metadata": {"title": "设备预测性维护解决方案建议书", "customer": "某制造集团"},
        })
        assert r.status_code == 200, f"❌ 导出请求失败 HTTP {r.status_code}: {r.text[:200]}"
        d = r.json()
        print(f"  POST /export/report -> {d.get('status')} wall={round(time.time()-t0,1)}s")

        assert d.get("status") == "completed", f"❌ 导出未完成: {d}"
        fname = d.get("file_name", "")
        assert fname.lower().endswith(".pptx"), f"❌ 文件名应为 .pptx，实际 {fname}"
        print(f"  ✅ 生成成功: {fname}")

        # 2. 下载
        r2 = await c.get(BASE + d["download_url"])
        assert r2.status_code == 200, f"❌ 下载失败 HTTP {r2.status_code}"
        data = r2.content
    print(f"  ✅ 下载 200，大小={len(data)} 字节")

    # 3. python-pptx 打开校验
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    n_slides = len(prs.slides)
    print(f"  ✅ 幻灯片数={n_slides}")
    assert n_slides >= 2, f"❌ 幻灯片数应 >= 2（封面+章节），实际 {n_slides}"

    cover_text = " ".join(s.text for s in prs.slides[0].shapes if hasattr(s, "text"))
    assert "方案" in cover_text or "华为云" in cover_text, f"❌ 封面标题异常: {cover_text[:80]}"
    print(f"  ✅ 封面内容: {cover_text[:80]}...")

    print("\nP2-D4 PPTX 导出验证全部通过 ✅")


if __name__ == "__main__":
    asyncio.run(main())
