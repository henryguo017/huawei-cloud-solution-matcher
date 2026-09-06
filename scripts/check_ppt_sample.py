# -*- coding: utf-8 -*-
"""v7 样张程序化自检：页数/越界/表格图表计数/开源技法抽查/金额去重/填实断言"""
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_FILL_TYPE

P = 'test_shots/huawei_red_ppt_sample.pptx'
prs = Presentation(P)
W, H = prs.slide_width, prs.slide_height
print('slides:', len(prs.slides._sldIdLst))

issues = []
n_tables = n_charts = 0
slides = list(prs.slides)
for idx, slide in enumerate(slides, 1):
    shapes = list(slide.shapes)
    for sh in shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            continue
        if l is None:
            continue
        if l < -10000 or t < -10000 or l + w > W + 10000 or t + h > H + 10000:
            issues.append('slide%d 越界: %s' % (idx, sh.shape_type))
        if t + h > H + 20000:
            issues.append('slide%d 底边出血: B=%.2f' % (idx, Emu(t + h).inches))
    for sh in shapes:
        if sh.has_table:
            n_tables += 1
        if getattr(sh, 'has_chart', False):
            n_charts += 1
    print('slide%-3d shapes=%d' % (idx, len(shapes)))

print('tables:', n_tables, ' charts:', n_charts)

# 技法抽查
s1 = slides[0]
ff = [sh for sh in s1.shapes if sh.shape_type is not None
      and 'FREEFORM' in str(sh.shape_type)]
t1 = '\n'.join(sh.text_frame.text for sh in s1.shapes if sh.has_text_frame)
# 字距 spc 属性抽查
spc_found = 0
for sh in s1.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                rPr = r._r.find(qn('a:rPr'))
                if rPr is not None and rPr.get('spc'):
                    spc_found += 1
print('s1 freeform(斜切几何, expect>=2):', len(ff),
      '| ghost SOLUTION:', ('SOLUTION' in t1), '| spc runs:', spc_found)

s2 = slides[1]
t2 = '\n'.join(sh.text_frame.text for sh in s2.shapes if sh.has_text_frame)
print('s2 ghost CONTENTS:', t2.count('CONTENTS'), '(expect>=2)')

# 每页幽灵章节号：抽 s3/s8
for sn in (3, 8):
    sl = slides[sn - 1]
    ts = [sh.text_frame.text for sh in sl.shapes if sh.has_text_frame]
    sec = ['01', '02', '03', '04', '05', '06', '07', '08'][sn - 3] if sn in (3, 8) else ''
    print('s%d ghost sec number %s:' % (sn, sec), sec in ts)

s12 = slides[11]
first = list(s12.shapes)[0]
try:
    print('s12 bg gradient:', first.fill.type == MSO_FILL_TYPE.GRADIENT)
except Exception as e:
    print('s12 bg gradient check err:', e)
t12 = '\n'.join(sh.text_frame.text for sh in s12.shapes if sh.has_text_frame)
print('s12 ghost THANKS:', ('THANKS' in t12),
      '| action chain:', all(k in t12 for k in ['确认方案范围', '商务细节洽谈', '12 周排期启动']))

s8 = slides[7]
t8 = [sh.table for sh in s8.shapes if sh.has_table]
hdr_fill = t8[0].cell(0, 0).fill.fore_color.rgb if t8 else None
print('s8 table header fill (expect E9EBEE 浅灰):', hdr_fill)

# ---- v9 去黑色量断言：全篇不得再有 INK(1A1C20) 填充块 ----
dark_fills = 0
for sl in slides:
    for sh in sl.shapes:
        try:
            if str(sh.fill.fore_color.rgb) == '1A1C20':
                dark_fills += 1
        except Exception:
            pass
print('全篇 INK 填充块 (expect 0):', dark_fills)
if dark_fills:
    issues.append('存在 INK 填充块')

# ---- v7 填实断言 ----
print('s2 方案摘要块:', '方案摘要' in t2 and 'EXECUTIVE SUMMARY' in t2)
t3 = '\n'.join(sh.text_frame.text for sh in slides[2].shapes if sh.has_text_frame)
print('s3 goals 验收行:', t3.count('验收：'), '(expect 4)')
t6 = '\n'.join(sh.text_frame.text for sh in slides[5].shapes if sh.has_text_frame)
print('s6 资源规格速览块:', '资源规格速览' in t6)

# ---- v7 金额去重断言（P8：总计唯一归宿=表格合计行）----
t8_text = '\n'.join(sh.text_frame.text for sh in s8.shapes if sh.has_text_frame)
tbl8_text = ''
for tb in t8:
    for r in range(len(tb.rows)):
        for c in range(len(tb.columns)):
            tbl8_text += tb.cell(r, c).text + '\n'
full8 = t8_text + '\n' + tbl8_text
money_ok = ('¥5,970' not in full8 and '¥60,889' not in full8
            and full8.count('60,889') == 1 and full8.count('5,969.50') == 1
            and '¥182,667' in full8 and full8.count('¥10,745') == 1)
print('s8 金额去重:', money_ok,
      '(¥5,970/¥60,889 消失 · 合计仅表格一处 · ¥182,667/¥10,745 各一次)')
if not money_ok:
    issues.append('s8 金额重复')

print('issues:', len(issues))
for i in issues:
    print(' -', i)
