# -*- coding: utf-8 -*-
"""空白审计：按 0.4 英寸水平带扫描每页，找出没有/极少形状覆盖的真空带。
目的：把"还有点空"从感觉变成坐标事实。排除全幅装饰（幽灵字/斜带/页眉页脚线）。
用法：python scripts/audit_whitespace.py [pptx路径]（缺省审计样张）"""
import sys

from pptx import Presentation
from pptx.util import Emu

P = sys.argv[1] if len(sys.argv) > 1 else 'test_shots/huawei_red_ppt_sample.pptx'
prs = Presentation(P)
W_IN = Emu(prs.slide_width).inches
H_IN = Emu(prs.slide_height).inches
BAND = 0.4          # 带高
Y0, Y1 = 1.2, 7.05  # 内容区（跳过页眉 KPI 行与页脚）

for idx, slide in enumerate(prs.slides, 1):
    boxes = []
    for sh in slide.shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            continue
        if l is None:
            continue
        x0, y0 = Emu(l).inches, Emu(t).inches
        x1, y1 = x0 + Emu(w).inches, y0 + Emu(h).inches
        # 排除装饰：全幅宽(>12.9)或超大盒(幽灵字/背景/图表框除外——图表算内容)
        if (x1 - x0) > 12.9:
            continue
        boxes.append((x0, y0, x1, y1))
    voids = []
    y = Y0
    while y < Y1:
        yb = min(y + BAND, Y1)
        # 与该带有实质相交(x 覆盖>0.6")的形状
        cov = 0.0
        for (x0, y0, x1, y1) in boxes:
            if y0 < yb and y1 > y:
                ov_x = min(x1, W_IN) - max(x0, 0)
                ov_y = min(y1, yb) - max(y0, y)
                if ov_x > 0.6 and ov_y > 0.12:
                    cov += ov_x
        ratio = cov / W_IN
        if ratio < 0.45:  # 覆盖不足 45% 视为真空带
            voids.append((round(y, 2), round(yb, 2), round(ratio, 2)))
        y = yb
    if voids:
        print('slide%-3d 真空带(y起-y止 覆盖率): %s' % (idx, voids))
print('audit done')
