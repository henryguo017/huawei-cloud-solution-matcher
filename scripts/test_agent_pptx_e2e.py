# -*- coding: utf-8 -*-
"""E2E：Agent 路径 PPTX 导出（无 cost_reference）——
_tool_generate_doc(fmt='pptx', content=终稿 markdown) → 引擎管线 12 页。
覆盖：意图修复后的真实用户场景"给我生成PPT"。"""
import asyncio
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.agent.tools import _tool_generate_doc

DRAFT = """# 智慧园区数字化上云解决方案建议书

## 一、项目背景
某省级制造集团下辖 3 个园区、入驻企业 380 余家、存量 IT 系统 8 套，资源分散、
运维成本高、数据孤岛问题突出。集团计划 2026 年启动整体上云，建设目标：
核心业务系统 100% 迁移、综合 IT 成本压降 25%、12 周内通过等保三级测评。

## 二、现状与核心痛点
P1 资源利用率不足 30%，扩容周期长达 6 周；P2 运维团队 12 人，故障平均恢复 4 小时；
P3 八套系统数据孤岛，跨系统对账依赖人工 Excel；P4 安全设备各自为政，
等保测评连续两年未通过。业务影响量化：年运维成本约 96 万元，故障损失月均 8 万元。

## 三、总体方案架构
采用华为云四层架构：接入层 WAF + ELB 统一入口；应用层 ECS c7.2xlarge.4 × 4 弹性承载；
数据层 RDS r7.2xlarge 主备高可用 + OBS 1TB 数据湖；安全层 DBSS 数据库审计、
HSS 主机安全、DSC 数据脱敏。设计要点：高可用、弹性伸缩、数据融合、安全合规、平滑迁移。

## 四、方案对比
对比自建 IDC：华为云 3 年 TCO 约 18.3 万元/月口径，自建约 26.9 万元，
自建另需一次性机房投入 300 万元；华为云扩容分钟级、SLA 99.95% 赔付，
综合成本优势约 32%。

## 五、成本估算
月费用结构：计算资源 640 元、数据库 1200 元、存储 400 元、网络带宽 300 元、
对象存储 120 元。年付 85 折，三年 TCO 直接成本约 182.7 万元。

## 六、实施路线
W1-2 调研与方案冻结；W3-6 资源开通与数据迁移；W7-10 应用割接与联调；
W11-12 等保测评与验收。里程碑 M1 方案冻结、M2 数据迁移完成、M3 割接完成、M4 验收。
"""


def main():
    obs = asyncio.run(_tool_generate_doc('pptx', content=DRAFT, report_type='solution'))
    data = json.loads(obs)
    print('status:', data.get('status'))
    print('file:', data.get('file_name'))
    assert data.get('status') == 'ok', data

    from pptx import Presentation
    p = Path(__file__).resolve().parents[1] / 'data' / 'exports' / data['file_name']
    prs = Presentation(str(p))
    slides = list(prs.slides)
    print('slides:', len(slides))
    assert len(slides) == 12, f'页数 {len(slides)} ≠ 12（可能降级 legacy）'

    # 成本页（第 8 页）应有合成明细表（无 cost_reference 场景）
    s8 = slides[7]
    tables = [sh for sh in s8.shapes if sh.has_table]
    assert tables, '成本页无明细表'
    for row in tables[0].table.rows:
        print(' | '.join(c.text for c in row.cells))
    print('AGENT-PATH E2E OK ->', p)


if __name__ == '__main__':
    main()
