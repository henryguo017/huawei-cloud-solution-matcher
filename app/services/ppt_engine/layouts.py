# -*- coding: utf-8 -*-
"""华为红售前 PPT 引擎 — 12 个满版版式（v9 样张拍板几何，禁改坐标）

每个版式 = render_<name>(prs, d)：几何/字体/色彩全部固化，内容只从 d 注入。
SLOT_SPEC 是每版式的槽位白名单 + 容量上限（P2 DeepSeek 填槽 schema 雏形，
engine.validate_deck 据此拒绝出稿）。

满版铁律：内容带 1.2~7.0" 铺满、块间距 ≤0.5"；改动任何坐标必须重跑
make→check→audit 三连且空白审计零整幅真空带。
"""
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

from . import tokens as T
from .primitives import (
    add_chart, add_grad_rect, add_oval, add_poly, add_rect, add_table,
    add_text, blank_slide, bullets, hline, page_header, section_head,
    set_cell, vline,
)


# ================================================================
# 槽位白名单 + 容量（engine.validate_deck 的出稿门禁依据）
# ================================================================
def _s(max_chars):
    return {'type': 'str', 'max_chars': max_chars}


def _l(max_items, item_chars=None):
    return {'type': 'list', 'max_items': max_items, 'item_chars': item_chars}


def _t(max_rows):
    """表格 dict 槽：max_rows = 表格行数硬上限（超出会被版式裁切，门禁拦截）"""
    return {'type': 'dict', 'max_rows': max_rows}


SLOT_SPEC = {
    'cover': {
        'title': _s(30), 'subtitle': _s(40),
        'cols': _l(3, 50), 'meta': _l(4, 24),
    },
    'toc': {
        'items': _l(8, 40), 'summary': {'type': 'dict'},
    },
    'background': {
        'header': {'type': 'header'}, 'paras': _l(3, 90),
        'kpis': _l(3, 24), 'chart': {'type': 'dict'},
        'path': {'type': 'dict'}, 'goals': {'type': 'dict'},
    },
    'pain_points': {
        'header': {'type': 'header'}, 'cards': _l(4, 60),
        'quant': _t(5), 'maps': {'type': 'dict'},
    },
    'requirements': {
        'header': {'type': 'header'}, 'cols': {'type': 'dict'},
        'resp': {'type': 'dict'}, 'matrix': {'type': 'dict'},
    },
    'architecture': {
        'header': {'type': 'header'}, 'layers': _l(3, 60),
        'points': {'type': 'dict'}, 'specs': {'type': 'dict'},
        'table': _t(4),
    },
    'compare': {
        'header': {'type': 'header'}, 'table': _t(8),
        'chart': {'type': 'dict'}, 'concls': {'type': 'dict'},
    },
    'cost': {
        'header': {'type': 'header'}, 'table': _t(7),
        'chart': {'type': 'dict'}, 'cards': _l(3, 40),
        'strategies': {'type': 'dict'},
    },
    'security': {
        'header': {'type': 'header'}, 'matrix': _t(6),
        'certs': _l(3, 40), 'plan': _t(4),
    },
    'roadmap': {
        'header': {'type': 'header'}, 'stages': _l(4, 60),
        'milestones': _t(4), 'risks': {'type': 'dict'},
    },
    'service': {
        'header': {'type': 'header'}, 'sla': _t(5),
        'services': {'type': 'dict'}, 'team': _t(4),
        'boundary': {'type': 'dict'},
    },
    'end': {'actions': _l(3, 12)},
}


# ================================================================
# 01 封面（大留白 + 编辑式三栏 + 斜切几何 + 幽灵字）
# ================================================================
def render_cover(prs, d):
    s = blank_slide(prs)
    add_poly(s, [(10.9, 0), (12.35, 0), (9.15, 3.3), (7.7, 3.3)], T.RED, alpha=6)
    add_poly(s, [(12.75, 0), (13.33, 0), (13.33, 2.5), (12.15, 2.5)], T.RED)
    add_text(s, Inches(5.6), Inches(1.02), Inches(7.73), Inches(1.3),
             [[("SOLUTION", 76, True, T.RED)]], align=PP_ALIGN.RIGHT,
             spc=600, alpha=5)
    # 品牌行（固定水印级内容）
    add_rect(s, Inches(1.0), Inches(0.58), Inches(0.2), Inches(0.2), T.RED)
    add_text(s, Inches(1.32), Inches(0.56), Inches(6), Inches(0.28),
             [[("华为云 HUAWEI CLOUD", 11, True, T.INK)]], spc=150)
    add_text(s, Inches(7.5), Inches(0.6), Inches(4.5), Inches(0.26),
             [[("SOLUTION PROPOSAL", 8.5, False, T.MUTE)]],
             align=PP_ALIGN.RIGHT, spc=300)
    # 主标题区
    add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.05),
             [[(d['title'], 42, True, T.INK)]])
    add_grad_rect(s, Inches(1.02), Inches(3.12), Inches(2.2), Inches(0.05),
                  angle=0)
    add_text(s, Inches(1.0), Inches(3.34), Inches(11.3), Inches(0.36),
             [[(d['subtitle'], 11.5, False, T.MUTE)]], spc=100)
    # 编辑式三栏
    cx = Inches(1.0)
    for t, txt in d['cols']:
        hline(s, cx, Inches(4.15), Inches(3.7))
        add_text(s, cx, Inches(4.3), Inches(3.7), Inches(0.32),
                 [[(t, 11.5, True, T.INK)]])
        add_text(s, cx, Inches(4.64), Inches(3.55), Inches(0.66),
                 [[(txt, 8.5, False, T.MUTE)]], line_spacing=1.2)
        cx += Inches(3.9)
    # 元信息行
    runs = []
    for k, v in d['meta']:
        runs += [("　　" if runs else "", 9.5, True, T.INK)]
        runs += [(k + "　", 9.5, True, T.INK), (v, 9.5, False, T.BODY)]
    add_text(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.3), [runs])
    # 页脚
    hline(s, Inches(1.0), Inches(6.85), Inches(11.333))
    add_text(s, Inches(1.0), Inches(6.98), Inches(8), Inches(0.26),
             [[(T.FOOTER_LEFT, 8, False, T.MUTE)]])
    add_text(s, Inches(9.5), Inches(6.98), Inches(2.833), Inches(0.26),
             [[(T.FOOTER_RIGHT, 8, False, T.MUTE)]], align=PP_ALIGN.RIGHT)


# ================================================================
# 02 目录（幽灵字 + 双列发丝线 + 方案摘要块）
# ================================================================
def render_toc(prs, d):
    s = blank_slide(prs)
    add_text(s, Inches(7.0), Inches(0.38), Inches(5.8), Inches(0.95),
             [[("CONTENTS", 44, True, T.RED)]], align=PP_ALIGN.RIGHT,
             spc=600, alpha=5)
    add_rect(s, Inches(1.0), Inches(1.05), Inches(0.55), Inches(0.07), T.RED)
    add_text(s, Inches(1.0), Inches(1.28), Inches(2.6), Inches(1.5), [
        [("目录", 26, True, T.INK)],
        [("CONTENTS", 10, False, T.MUTE)],
    ], space_after=8, spc=300)
    add_text(s, Inches(1.0), Inches(6.5), Inches(2.8), Inches(0.6),
             [[("华为云 · 智慧园区解决方案", 8.5, False, T.MUTE)],
              [("2026-09 · V1.0 · 商密", 8.5, False, T.MUTE)]], space_after=3)
    items = d['items']
    col_x = [Inches(4.0), Inches(8.65)]
    for i, (num, title, sub) in enumerate(items):
        cx = col_x[i // 4]
        cy = Inches(1.0 + (i % 4) * 1.52)
        add_text(s, cx, cy, Inches(0.7), Inches(0.55),
                 [[(num, 15, True, T.RED)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.72), cy - Inches(0.02), Inches(3.5),
                 Inches(0.4), [[(title, 12.5, True, T.INK)]])
        add_text(s, cx + Inches(0.72), cy + Inches(0.36), Inches(3.5),
                 Inches(0.35), [[(sub, 8, False, T.MUTE)]])
        hline(s, cx, cy + Inches(0.88), Inches(4.22))
    # 左栏方案摘要块
    sm = d['summary']
    hline(s, Inches(1.0), Inches(3.15), Inches(2.8))
    add_text(s, Inches(1.0), Inches(3.3), Inches(2.8), Inches(0.3),
             [[("方案摘要", 11.5, True, T.INK)]])
    add_text(s, Inches(1.0), Inches(3.58), Inches(2.8), Inches(0.24),
             [[("EXECUTIVE SUMMARY", 7.5, False, T.MUTE)]], spc=200)
    bullets(s, Inches(1.0), Inches(3.94), Inches(2.85),
            sm['items'], size=8, gap=4, dot=False)
    add_text(s, Inches(1.0), Inches(5.42), Inches(2.85), Inches(0.9),
             [[("结论　", 8.5, True, T.RED), (sm['conclusion'], 8.5, False,
                                               T.BODY)]],
             line_spacing=1.25)
    hline(s, Inches(1.0), Inches(7.075), Inches(11.333))
    add_text(s, Inches(1.0), Inches(7.14), Inches(8), Inches(0.26),
             [[(T.FOOTER_LEFT, 7.5, False, T.MUTE)]])
    add_text(s, Inches(11.6), Inches(7.14), Inches(0.733), Inches(0.26),
             [[("02 / %d" % T.TOTAL_PAGES, 8, False, T.MUTE)]],
             align=PP_ALIGN.RIGHT)


# ================================================================
# 03 背景与目标（论述 + KPI + 折线图 + 路径时间轴 + 4 目标）
# ================================================================
def render_background(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    section_head(s, T.MX, Inches(1.32), Inches(6.7), "项目背景")
    bullets(s, T.MX + Inches(0.05), Inches(1.82), Inches(6.6),
            d['paras'], size=9.5, gap=10, dot=False, text_color=T.BODY)
    # 路径时间轴（左栏下部）
    p = d['path']
    hline(s, T.MX, Inches(3.5), Inches(6.7))
    add_text(s, T.MX, Inches(3.62), Inches(4), Inches(0.28),
             [[(p['title'], 11.5, True, T.INK)]])
    add_grad_rect(s, Inches(0.62), Inches(4.28), Inches(6.2), Inches(0.03),
                  angle=0)
    px = 0.75
    nodes = p['nodes']
    for i, (yr, ev) in enumerate(nodes):
        last = (i == len(nodes) - 1)
        add_oval(s, Inches(px - 0.055), Inches(4.265), Inches(0.11),
                 Inches(0.11), T.RED)
        add_text(s, Inches(px - 0.55), Inches(3.94), Inches(1.1), Inches(0.26),
                 [[(yr, 9.5, True, T.RED if last else T.INK)]],
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(px - 0.55), Inches(4.44), Inches(1.1), Inches(0.24),
                 [[(ev, 8, False, T.RED if last else T.BODY)]],
                 align=PP_ALIGN.CENTER)
        px += 1.45
    # 右：KPI + 折线图
    kx = Inches(7.4)
    for num, sub in d['kpis']:
        hline(s, kx, Inches(1.36), Inches(1.7))
        add_text(s, kx, Inches(1.48), Inches(1.7), Inches(0.36),
                 [[(num, 15, True, T.RED)]])
        add_text(s, kx, Inches(1.84), Inches(1.7), Inches(0.34),
                 [[(sub, 7, False, T.MUTE)]], line_spacing=1.05)
        kx += Inches(1.82)
    c = d['chart']
    section_head(s, Inches(7.4), Inches(2.4), Inches(5.48), c['title'])
    add_chart(s, Inches(7.4), Inches(2.84), Inches(5.48), Inches(1.9),
              XL_CHART_TYPE.LINE_MARKERS, c['cats'], c['vals'], c['name'],
              num_fmt=c.get('num_fmt', '0'))
    # 下半：4 目标（含验收行）
    g = d['goals']
    section_head(s, T.MX, Inches(4.92), T.CW, g['title'])
    gx = T.MX
    for t, m1, m2, m3, acc in g['rows']:
        hline(s, gx, Inches(5.4), Inches(3.02))
        add_text(s, gx, Inches(5.5), Inches(2.9), Inches(0.3),
                 [[(t, 10.5, True, T.INK)]])
        for i, m in enumerate([m1, m2]):
            add_text(s, gx, Inches(5.8 + i * 0.27), Inches(2.9), Inches(0.26),
                     [[("· ", 8.5, True, T.RED), (m, 8.5, False, T.BODY)]])
        add_text(s, gx, Inches(6.36), Inches(2.9), Inches(0.24),
                 [[(m3, 8.5, True, T.RED)]])
        add_text(s, gx, Inches(6.62), Inches(2.9), Inches(0.24),
                 [[(acc, 7.5, False, T.MUTE)]])
        gx += Inches(3.14)


# ================================================================
# 04 痛点（2×2 卡片含 hero 数字 + 量化表 + 对策映射）
# ================================================================
def render_pain_points(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    cw, ch = Inches(3.82), Inches(2.6)
    for i, (t, cur, impact, hnum, hlabel) in enumerate(d['cards']):
        cx = T.MX + (cw + Inches(0.18)) * (i % 2)
        cy = Inches(1.32) + (ch + Inches(0.16)) * (i // 2)
        hline(s, cx, cy, cw)
        add_text(s, cx, cy + Inches(0.12), cw - Inches(0.3), Inches(0.3),
                 [[(t, 11, True, T.INK)]])
        bullets(s, cx + Inches(0.02), cy + Inches(0.56), cw - Inches(0.1),
                [("现状　", cur), ("影响　", impact)], size=9, gap=11)
        add_text(s, cx, cy + Inches(1.95), cw - Inches(0.3), Inches(0.4),
                 [[(hnum, 16, True, T.RED),
                   ("　" + hlabel, 7.5, False, T.MUTE)]])
    # 右：量化表 + 映射
    q = d['quant']
    tx = Inches(8.5)
    section_head(s, tx, Inches(1.32), Inches(4.38), q['title'])
    add_table(s, tx, Inches(1.78), Inches(4.38), q['headers'], q['rows'],
              q.get('widths'), row_h=0.33,
              col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])
    m = d['maps']
    section_head(s, tx, Inches(3.98), Inches(4.38), m['title'])
    my = Inches(4.42)
    for pid, act in m['pairs']:
        add_rect(s, tx, my + Inches(0.08), Inches(0.36), Inches(0.36), T.RED)
        add_text(s, tx, my + Inches(0.08), Inches(0.36), Inches(0.36),
                 [[(pid, 9, True, T.WHITE)]], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + Inches(0.52), my, Inches(3.86), Inches(0.56),
                 [[(act, 9, False, T.BODY)]], anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.1)
        my += Inches(0.62)


# ================================================================
# 05 需求（三栏 + 本期响应 + 优先级矩阵）
# ================================================================
def render_requirements(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    colw = Inches(4.02)
    r = d['resp']
    for ci, (title, items) in enumerate(d['cols']):
        x = T.MX + (colw + Inches(0.18)) * ci
        add_rect(s, x, Inches(1.38), Inches(0.09), Inches(0.09), T.RED)
        add_text(s, x + Inches(0.2), Inches(1.32), colw - Inches(0.25),
                 Inches(0.28), [[(title, 11.5, True, T.INK)]])
        hline(s, x, Inches(1.66), colw)
        bullets(s, x + Inches(0.02), Inches(1.84), colw - Inches(0.12),
                items, size=9, gap=7, dot=False, text_color=T.BODY)
        hline(s, x, Inches(3.72), colw)
        add_text(s, x + Inches(0.02), Inches(3.84), colw - Inches(0.1),
                 Inches(0.26), [[(r['title'], 9.5, True, T.RED)]])
        bullets(s, x + Inches(0.02), Inches(4.14), colw - Inches(0.12),
                r['cols'][ci], size=8.5, gap=5, dot=False)
    vline(s, Inches(4.56), Inches(1.36), Inches(3.5))
    vline(s, Inches(8.76), Inches(1.36), Inches(3.5))
    # 优先级矩阵
    m = d['matrix']
    section_head(s, T.MX, Inches(4.92), T.CW, m['title'])
    style_map = {'grad': (None, T.WHITE, None),
                 'outline_red': (T.WHITE, T.RED, T.RED),
                 'outline_gray': (T.WHITE, T.MUTE, T.RULE)}
    gx = T.MX
    for t, kind, desc in m['groups']:
        fill, tcol, border = style_map[kind]
        if kind == 'grad':
            add_grad_rect(s, gx, Inches(5.38), Inches(4.02), Inches(1.26),
                          angle=90)
        else:
            add_rect(s, gx, Inches(5.38), Inches(4.02), Inches(1.26), fill,
                     line_color=border, line_w=1.0)
        add_text(s, gx + Inches(0.18), Inches(5.5), Inches(3.7), Inches(0.28),
                 [[(t, 11, True, tcol)]])
        add_text(s, gx + Inches(0.18), Inches(5.84), Inches(3.72),
                 Inches(0.7),
                 [[(desc, 9, False, T.WHITE if kind == 'grad' else T.MUTE)]],
                 line_spacing=1.3)
        gx += Inches(4.2)


# ================================================================
# 06 架构（浅红命名块 + 设计要点 + 规格速览 + 选型表）
# ================================================================
def render_architecture(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    y = Inches(1.32)
    for name, tag, comps in d['layers']:
        hh = Inches(1.16)
        add_rect(s, T.MX, y, Inches(8.0), hh, T.WHITE, line_color=T.RULE,
                 line_w=1.0)
        add_rect(s, T.MX, y, Inches(0.05), hh, T.RED)
        add_rect(s, T.MX + Inches(0.05), y, Inches(1.3), hh, T.RED_SOFT)
        add_text(s, T.MX + Inches(0.05), y + Inches(0.24), Inches(1.3),
                 Inches(0.38), [[(name, 12.5, True, T.RED)]],
                 align=PP_ALIGN.CENTER)
        add_text(s, T.MX + Inches(0.05), y + Inches(0.64), Inches(1.3),
                 Inches(0.28), [[(tag, 7.5, False, T.MUTE)]],
                 align=PP_ALIGN.CENTER)
        cx = Inches(2.1)
        for comp, wd in comps:
            add_rect(s, cx, y + Inches(0.3), Inches(wd), Inches(0.56),
                     T.WHITE, line_color=T.RULE, line_w=1.0)
            add_text(s, cx, y + Inches(0.3), Inches(wd), Inches(0.56),
                     [[(comp, 8.5, True, T.INK)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += Inches(wd + 0.13)
        y += hh + Inches(0.12)
    # 右侧设计要点 + 规格速览
    p = d['points']
    section_head(s, Inches(8.62), Inches(1.32), Inches(4.26), p['title'])
    bullets(s, Inches(8.66), Inches(1.84), Inches(4.24),
            p['rows'], size=9, gap=4.5)
    sp = d['specs']
    hline(s, Inches(8.62), Inches(3.24), Inches(4.26))
    add_text(s, Inches(8.62), Inches(3.38), Inches(4.26), Inches(0.28),
             [[(sp['title'], 11.5, True, T.INK)]])
    bullets(s, Inches(8.66), Inches(3.74), Inches(4.24),
            sp['rows'], size=8.5, gap=3.5)
    # 底部选型依据表
    tb = d['table']
    section_head(s, T.MX, Inches(5.02), T.CW, tb['title'])
    add_table(s, T.MX, Inches(5.48), T.CW, tb['headers'], tb['rows'],
              tb.get('widths'), row_h=0.32, font_size=8.5)


# ================================================================
# 07 对比（8 行表 + TCO 柱图 + 结论）
# ================================================================
def render_compare(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    tb = d['table']
    add_table(s, T.MX, Inches(1.32), Inches(7.6), tb['headers'], tb['rows'],
              tb.get('widths'), row_h=0.31, header_h=0.34,
              col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
                          PP_ALIGN.CENTER],
              font_size=8)
    c = d['chart']
    section_head(s, T.MX, Inches(4.28), Inches(7.6), c['title'])
    add_chart(s, Inches(0.55), Inches(4.76), Inches(7.3), Inches(2.2),
              XL_CHART_TYPE.COLUMN_CLUSTERED, c['cats'], c['vals'],
              c['name'], num_fmt=c.get('num_fmt', '0.0'),
              point_colors=[T.RED if x == 'red' else T.GRAY_BAR
                            for x in c['colors']])
    cc = d['concls']
    section_head(s, Inches(8.35), Inches(1.32), Inches(4.53), cc['title'])
    cy = Inches(1.82)
    for t, txt in cc['rows']:
        hline(s, Inches(8.35), cy, Inches(4.53))
        add_text(s, Inches(8.35), cy + Inches(0.12), Inches(4.4), Inches(0.3),
                 [[(t, 11, True, T.RED)]])
        add_text(s, Inches(8.35), cy + Inches(0.44), Inches(4.4), Inches(0.72),
                 [[(txt, 8.5, False, T.BODY)]], line_spacing=1.18)
        cy += Inches(1.34)


# ================================================================
# 08 成本（明细表含合计行 + 柱图 + 要点卡 + 优化策略）
# ================================================================
def render_cost(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    tb = d.get('table')
    if tb:
        add_table(s, T.MX, Inches(1.32), T.CW, tb['headers'], tb['rows'],
                  tb.get('widths'), row_h=0.32, header_h=0.34,
                  col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                              PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
                              PP_ALIGN.RIGHT, PP_ALIGN.CENTER],
                  font_size=8.5)
        # 合计行强调
        tbl = [sh for sh in s.shapes if sh.has_table][0].table
        tri = tb.get('total_row_idx', len(tb['rows']))
        accent = tuple(tb.get('accent_cols', []))
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                  PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
                  PP_ALIGN.RIGHT, PP_ALIGN.CENTER]
        for j in range(len(tb['headers'])):
            set_cell(tbl.cell(tri + 1, j), tb['rows'][tri][j], size=8.5,
                     bold=True,
                     color=T.RED if j in accent else T.INK, fill=T.RED_SOFT,
                     align=aligns[j])
    # 下左：柱图（WAF 标红）
    c = d['chart']
    section_head(s, T.MX, Inches(3.86), Inches(5.9), c['title'])
    add_chart(s, Inches(0.55), Inches(4.34), Inches(5.7), Inches(2.45),
              XL_CHART_TYPE.COLUMN_CLUSTERED, c['cats'], c['vals'],
              c['name'], num_fmt=c.get('num_fmt', '#,##0'),
              point_colors=[T.RED if x == 'red' else T.GRAY_BAR
                            for x in c['colors']])
    # 下中：要点卡（金额去重：只放派生/结构信息）
    kx = Inches(6.55)
    ky = Inches(3.86)
    for t, num, sub, style in d['cards']:
        col = T.RED if style == 'red' else T.INK
        hline(s, kx, ky, Inches(3.0))
        add_text(s, kx, ky + Inches(0.1), Inches(2.85), Inches(0.24),
                 [[(t, 8, False, T.MUTE)]])
        add_text(s, kx, ky + Inches(0.32), Inches(2.85), Inches(0.36),
                 [[(num, 16, True, col)]])
        add_text(s, kx, ky + Inches(0.62), Inches(2.85), Inches(0.24),
                 [[(sub, 7.5, False, T.MUTE)]])
        ky += Inches(0.94)
    # 下右：优化策略
    st = d['strategies']
    section_head(s, Inches(9.75), Inches(3.86), Inches(3.13), st['title'])
    bullets(s, Inches(9.78), Inches(4.34), Inches(3.1),
            st['rows'], size=8.5, gap=6)


# ================================================================
# 09 安全合规（等保矩阵 + 认证 + 测评计划）
# ================================================================
def render_security(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    m = d['matrix']
    add_table(s, T.MX, Inches(1.32), T.CW, m['headers'], m['rows'],
              m.get('widths'), row_h=0.36, header_h=0.34,
              col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
                          PP_ALIGN.CENTER],
              font_size=8.5)
    cx = T.MX
    for t, txt in d['certs']:
        hline(s, cx, Inches(4.0), Inches(4.04))
        add_text(s, cx, Inches(4.14), Inches(3.9), Inches(0.3),
                 [[(t, 10.5, True, T.INK)]])
        add_text(s, cx, Inches(4.46), Inches(3.9), Inches(0.4),
                 [[(txt, 8.5, False, T.MUTE)]], line_spacing=1.1)
        cx += Inches(4.2)
    p = d['plan']
    section_head(s, T.MX, Inches(5.02), T.CW, p['title'])
    add_table(s, T.MX, Inches(5.48), T.CW, p['headers'], p['rows'],
              p.get('widths'), row_h=0.32, font_size=8.5)


# ================================================================
# 10 实施路线（时间轴 + 里程碑 + 风险）
# ================================================================
def render_roadmap(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    add_grad_rect(s, Inches(0.95), Inches(1.78), Inches(11.45), Inches(0.035),
                  angle=0)
    x0, step = 1.5, 2.98
    for i, (wk, title, delivs) in enumerate(d['stages']):
        cx = x0 + i * step
        add_oval(s, Inches(cx - 0.07), Inches(1.73), Inches(0.14), Inches(0.14),
                 T.RED)
        add_text(s, Inches(cx - 0.8), Inches(1.3), Inches(1.6), Inches(0.26),
                 [[(wk, 9.5, True, T.RED)]], align=PP_ALIGN.CENTER)
        add_text(s, Inches(cx - 1.05), Inches(2.0), Inches(2.1), Inches(0.35),
                 [[(title, 12.5, True, T.INK)]], align=PP_ALIGN.CENTER)
        add_text(s, Inches(cx - 1.3), Inches(2.38), Inches(2.6), Inches(0.55),
                 [[("交付物：", 7.8, True, T.MUTE),
                   (" · ".join(delivs), 7.8, False, T.MUTE)]],
                 align=PP_ALIGN.CENTER, line_spacing=1.1)
    m = d['milestones']
    section_head(s, T.MX, Inches(3.02), T.CW, m['title'])
    add_table(s, T.MX, Inches(3.48), T.CW, m['headers'], m['rows'],
              m.get('widths'), row_h=0.33, font_size=8.5)
    r = d['risks']
    section_head(s, T.MX, Inches(5.3), T.CW, r['title'])
    rx = T.MX
    for t, lvl, desc, plan in r['rows']:
        hline(s, rx, Inches(5.74), Inches(4.04))
        add_text(s, rx, Inches(5.86), Inches(3.9), Inches(0.28),
                 [[(t, 10, True, T.INK),
                   ("　风险 " + lvl, 8.5, True, T.RED if lvl == "高"
                    else T.MUTE)]])
        add_text(s, rx, Inches(6.16), Inches(3.9), Inches(0.28),
                 [[(desc, 8, False, T.MUTE)]])
        add_text(s, rx, Inches(6.46), Inches(3.9), Inches(0.45),
                 [[("应对：", 8, True, T.INK), (plan, 8, False, T.BODY)]],
                 line_spacing=1.12)
        rx += Inches(4.2)


# ================================================================
# 11 服务保障（SLA + 服务清单 + 团队 + 边界）
# ================================================================
def render_service(prs, d):
    s = blank_slide(prs)
    h = d['header']
    page_header(s, h['sec'], h['sec_name'], h['headline'], h['page_no'],
                stats=h.get('stats'), note=h.get('note'))
    sla = d['sla']
    section_head(s, T.MX, Inches(1.32), Inches(7.55), sla['title'])
    add_table(s, T.MX, Inches(1.78), Inches(7.55), sla['headers'],
              sla['rows'], sla.get('widths'), row_h=0.34, font_size=8.5)
    sv = d['services']
    section_head(s, Inches(8.25), Inches(1.32), Inches(4.63), sv['title'])
    bullets(s, Inches(8.3), Inches(1.84), Inches(4.55),
            sv['rows'], size=9, gap=4.5)
    tm = d['team']
    section_head(s, T.MX, Inches(3.92), T.CW, tm['title'])
    add_table(s, T.MX, Inches(4.38), T.CW, tm['headers'], tm['rows'],
              tm.get('widths'), row_h=0.33, font_size=8.5)
    b = d['boundary']
    section_head(s, T.MX, Inches(6.0), T.CW, b['title'])
    bullets(s, T.MX + Inches(0.05), Inches(6.48), T.CW - Inches(0.1),
            b['rows'], size=8.5, gap=2.5, dot=False, text_color=T.MUTE)


# ================================================================
# 12 结尾（红底大字 + 行动链）
# ================================================================
def render_end(prs, d):
    s = blank_slide(prs)
    add_grad_rect(s, Inches(0), Inches(0), T.SLIDE_W, T.SLIDE_H,
                  T.RED, T.RED_DEEP, angle=45)
    add_rect(s, Inches(0), Inches(0), Inches(0.12), T.SLIDE_H, T.RED_DEEP)
    add_text(s, Inches(0), Inches(0.72), T.SLIDE_W, Inches(1.5),
             [[("THANKS", 88, True, T.WHITE)]], align=PP_ALIGN.CENTER,
             spc=800, alpha=8)
    add_text(s, Inches(0), Inches(2.42), T.SLIDE_W, Inches(1.2),
             [[("谢 谢 观 看", 40, True, T.WHITE)]], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.7), T.SLIDE_W, Inches(0.4),
             [[("THANKS FOR WATCHING", 11, False, T.TAG)]],
             align=PP_ALIGN.CENTER, spc=400)
    add_rect(s, Inches(6.42), Inches(4.26), Inches(0.5), Inches(0.035),
             T.WHITE)
    acts = d['actions']
    runs = [("下一步　", 11, True, T.TAG)]
    for i, a in enumerate(acts):
        runs.append((a, 13, True, T.WHITE))
        if i < len(acts) - 1:
            runs.append(("　→　", 12, True, T.TAG))
    add_text(s, Inches(0), Inches(4.62), T.SLIDE_W, Inches(0.4), [runs],
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(6.3), T.SLIDE_W, Inches(0.4),
             [[("华为云 · 解决方案智能匹配自动生成 ｜ " + T.FOOTER_RIGHT,
                9.5, False, T.TAG)]], align=PP_ALIGN.CENTER)


LAYOUTS = {
    'cover': render_cover,
    'toc': render_toc,
    'background': render_background,
    'pain_points': render_pain_points,
    'requirements': render_requirements,
    'architecture': render_architecture,
    'compare': render_compare,
    'cost': render_cost,
    'security': render_security,
    'roadmap': render_roadmap,
    'service': render_service,
    'end': render_end,
}
