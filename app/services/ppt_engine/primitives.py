# -*- coding: utf-8 -*-
"""华为红售前 PPT 引擎 — 绘制原语（几何/字体/字距/透明度全部在此锁定）"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from . import tokens as T

MIDDLE = MSO_ANCHOR.MIDDLE
CENTER = PP_ALIGN.CENTER
LEFT = PP_ALIGN.LEFT
RIGHT = PP_ALIGN.RIGHT


def _set_run(run, text, size, bold, color, spc=None, alpha=None):
    """spc: 字距（1/100 pt，如 300=3pt）；alpha: 文字不透明度百分比（如 5=5%）"""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Arial'          # 西文/数字用 Arial，中文走 ea 字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', T.FONT)
    if spc is not None:
        rPr.set('spc', str(int(spc)))
    if alpha is not None:
        sf = rPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = clr.makeelement(qn('a:alpha'),
                                    {'val': str(int(alpha * 1000))})
                clr.append(a)


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space_after=3, line_spacing=1.0,
             spc=None, alpha=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, s, b, c) in para:
            _set_run(p.add_run(), t, s, b, c, spc=spc, alpha=alpha)
    return tb


def add_rect(slide, x, y, w, h, fill, rounded=False, radius=0.08,
             line_color=None, line_w=0.75):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, w, h, fill, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def hline(slide, x, y, w, color=None, thick=0.011):
    add_rect(slide, x, y, w, Inches(thick), color or T.RULE)


def vline(slide, x, y, h, color=None, thick=0.011):
    add_rect(slide, x, y, Inches(thick), h, color or T.RULE)


def set_shape_alpha(shp, pct):
    """形状填充不透明度（pct 百分比，如 6 → alpha 6%）"""
    spPr = shp._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        return
    clr = sf.find(qn('a:srgbClr'))
    if clr is None:
        return
    a = clr.makeelement(qn('a:alpha'), {'val': str(int(pct * 1000))})
    clr.append(a)


def add_grad_rect(slide, x, y, w, h, c1=None, c2=None, angle=90):
    """单色系渐变矩形（品牌红→深红，角度：0=向右 90=向下 45=对角）"""
    c1 = c1 or T.RED
    c2 = c2 or T.RED_DEEP
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.fill.gradient()
        stops = shp.fill.gradient_stops
        stops[0].color.rgb = c1
        stops[0].position = 0.0
        stops[1].color.rgb = c2
        stops[1].position = 1.0
        try:
            shp.fill.gradient_angle = angle
        except Exception:
            pass
    except Exception:
        shp.fill.solid()
        shp.fill.fore_color.rgb = c1
    return shp


def add_poly(slide, pts, fill=None, alpha=None):
    """自由路径多边形（SVG path 思维）；pts: [(x_inch, y_inch), ...] 顺时针"""
    fill = fill or T.RED
    fb = slide.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]))
    fb.add_line_segments([(Inches(px), Inches(py)) for px, py in pts[1:]],
                         close=True)
    shp = fb.convert_to_shape()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if alpha is not None:
        set_shape_alpha(shp, alpha)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def set_cell(cell, text, size=8.5, bold=False, color=None, fill=None,
             align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.015)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill is not None else T.WHITE
    tf = cell.text_frame
    tf.clear()  # 防重复写入：合计行会被二次 set_cell 强调，不清空会双倍显示金额
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size, bold, color or T.BODY)


def add_table(slide, x, y, w, headers, rows, col_widths=None, row_h=0.32,
              header_h=0.34, col_aligns=None, font_size=8.5):
    gf = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w,
                                Inches(header_h + row_h * len(rows)))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    if not col_widths:
        # 未指定列宽时按列数均分表宽（LLM 数据无 widths 槽时的兜底）
        eq = Emu(w).inches / len(headers)
        col_widths = [eq] * len(headers)
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)
    tbl.rows[0].height = Inches(header_h)
    for j, htxt in enumerate(headers):
        set_cell(tbl.cell(0, j), htxt, size=font_size, bold=True,
                 color=T.INK, fill=T.THEAD, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        tbl.rows[r].height = Inches(row_h)
        band = T.WHITE if r % 2 == 1 else T.FAINT
        for j, val in enumerate(row):
            align = (col_aligns[j] if col_aligns else PP_ALIGN.LEFT)
            set_cell(tbl.cell(r, j), str(val), size=font_size,
                     bold=(j == 0), color=T.INK if j == 0 else T.BODY,
                     fill=band, align=align)
    return gf


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def stat_row(slide, stats, right=12.883, y=0.54, w=1.16, gap=0.12):
    """右上 KPI 对：红色数字 + 灰色标签，发丝竖线分隔，无边框"""
    total = len(stats) * w + (len(stats) - 1) * gap
    x = right - total
    for i, (num, label) in enumerate(stats):
        cx = Inches(x + i * (w + gap))
        add_text(slide, cx, Inches(y), Inches(w), Inches(0.26),
                 [[(num, 12.5, True, T.RED)]], align=PP_ALIGN.CENTER)
        add_text(slide, cx, Inches(y + 0.27), Inches(w), Inches(0.2),
                 [[(label, 7.5, False, T.MUTE)]], align=PP_ALIGN.CENTER)
        if i > 0:
            vline(slide, Inches(x + i * (w + gap) - gap / 2),
                  Inches(y + 0.02), Inches(0.42))
    return x


def page_header(slide, sec, sec_name, headline, page_no, stats=None, note=None):
    """页眉：幽灵章节号 + 红色 kicker + 大标题 / 右上 KPI 对；页脚：来源 + 页码"""
    # 幽灵章节号（5% 红，背景层次）
    add_text(slide, Inches(10.9), Inches(0.08), Inches(1.98), Inches(0.9),
             [[(sec, 54, True, T.RED)]], align=PP_ALIGN.RIGHT, spc=200, alpha=5)
    add_text(slide, T.MX, Inches(0.28), Inches(6), Inches(0.22),
             [[("%s · %s" % (sec, sec_name), 9, True, T.RED)]], spc=100)
    add_text(slide, Inches(7.0), Inches(0.28), Inches(5.883), Inches(0.22),
             [[("华为云 HUAWEI CLOUD", 8, False, T.MUTE)]], align=PP_ALIGN.RIGHT,
             spc=150)
    add_text(slide, T.MX, Inches(0.5), Inches(8.5), Inches(0.46),
             [[(headline, 22, True, T.INK)]], anchor=MSO_ANCHOR.MIDDLE)
    if stats:
        stat_row(slide, stats)
    hline(slide, T.MX, Inches(1.16), T.CW)
    add_grad_rect(slide, T.MX, Inches(1.14), Inches(0.9), Inches(0.05), angle=0)
    hline(slide, T.MX, Inches(7.075), T.CW)
    if note:
        add_text(slide, T.MX, Inches(7.12), Inches(10.9), Inches(0.26),
                 [[("说明：", 7.5, True, T.MUTE), (note, 7.5, False, T.MUTE)]])
    add_text(slide, Inches(11.6), Inches(7.12), Inches(1.283), Inches(0.26),
             [[("%s / %d" % (page_no, T.TOTAL_PAGES), 8, False, T.MUTE)]],
             align=PP_ALIGN.RIGHT)


def section_head(slide, x, y, w, title):
    """小节头：红方块 + 标题 + 发丝线（替代色带填底）"""
    add_rect(slide, x, Inches(Emu(y).inches + 0.06), Inches(0.09),
             Inches(0.09), T.RED)
    add_text(slide, x + Inches(0.2), y, w - Inches(0.25), Inches(0.26),
             [[(title, 11.5, True, T.INK)]])
    hline(slide, x, y + Inches(0.32), w)


def rule_card(slide, x, y, w, h):
    """编辑式卡片：仅顶部发丝线，无底色无边框"""
    hline(slide, x, y, w)


def bullets(slide, x, y, w, items, size=9, gap=3.5, dot=True,
            label_color=None, text_color=None):
    """紧凑列表；items: [(label, desc)] 或 str；返回结束 y(Inches float)"""
    label_color = label_color or T.INK
    text_color = text_color or T.BODY
    cy = Emu(y).inches
    for it in items:
        # label+desc 对：兼容 tuple 与 list（LLM JSON 产物一律是 list，样张是 tuple）
        if (isinstance(it, (tuple, list)) and len(it) == 2
                and all(isinstance(p, str) for p in it)):
            label, desc = it
        else:
            label, desc = None, it
        desc = desc if isinstance(desc, str) else str(desc)
        if dot:
            add_rect(slide, x, Inches(cy + 0.05), Inches(0.055),
                     Inches(0.055), T.RED)
        runs = []
        if label:
            runs.append((label, size, True, label_color))
        runs.append((desc, size, False, text_color))
        add_text(slide, x + Inches(0.14), Inches(cy), w - Inches(0.14),
                 Inches(0.5), [runs], line_spacing=1.08)
        chars_per_line = max(8, int(Emu(w - Inches(0.14)).inches * 11.8))
        est_lines = max(1, -(-len((label or '') + desc) // chars_per_line))
        cy += est_lines * (size / 72.0) * 1.22 + gap / 72.0
    return cy


def add_chart(slide, x, y, w, h, ctype, categories, values, series_name,
              num_fmt='#,##0', color=None, label_size=7.5, axis_size=8,
              point_colors=None):
    """原生图表；point_colors: 按数据点着色（柱状图红/灰语义）"""
    color = color or T.RED
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series(series_name, values, num_fmt)
    gf = slide.shapes.add_chart(ctype, x, y, w, h, cd)
    ch = gf.chart
    ch.has_legend = False
    ch.font.size = Pt(axis_size)
    ch.font.name = T.FONT
    plot = ch.plots[0]
    plot.gap_width = 55
    ser = plot.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = color
    ser.format.line.fill.background()
    if point_colors:
        try:
            for i, pc in enumerate(point_colors):
                pt = ser.points[i]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = pc
        except Exception:
            pass
    if ctype == XL_CHART_TYPE.LINE_MARKERS:
        try:
            from pptx.enum.chart import XL_MARKER_STYLE
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(1.75)
            ser.smooth = False
            ser.marker.style = XL_MARKER_STYLE.CIRCLE
            ser.marker.size = 5
            ser.marker.format.fill.solid()
            ser.marker.format.fill.fore_color.rgb = color
            ser.marker.format.line.color.rgb = T.WHITE
        except Exception:
            pass
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(label_size)
    dl.font.color.rgb = T.BODY
    dl.font.bold = True
    dl.number_format_is_linked = False
    dl.number_format = num_fmt
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(axis_size)
    ca.format.line.color.rgb = T.RULE
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = T.FAINT
    va.tick_labels.font.size = Pt(7)
    return ch
