import time
import re
from datetime import datetime
from typing import Dict, Any, List
from pathlib import Path

from app.models.export_models import (
    ExportFormat, TaskStatus, ReportType, 
    ExportTask
)
from app.utils.word_generator import WordGenerator


# 模块级单例：保证 Agent 的 generate_doc 工具与 /api/export/download 共用同一任务表，
# 使 generate_report 返回的 task_id 在下载路由中可被查到（否则跨实例 → 404）。
_report_generator_instance = None


def get_report_generator() -> "ReportGeneratorService":
    """获取 ReportGeneratorService 单例（延迟实例化）。"""
    global _report_generator_instance
    if _report_generator_instance is None:
        _report_generator_instance = ReportGeneratorService()
    return _report_generator_instance


def _fmt_num(n) -> str:
    """把数字格式化为千分位字符串（整数不带小数，小数最多两位）。"""
    try:
        v = float(n)
    except (TypeError, ValueError):
        return str(n)
    if v == int(v):
        return f"{int(v):,}"
    return f"{round(v, 2):,}"


class ReportGeneratorService:
    """报告生成服务"""
    
    def __init__(self):
        self.export_dir = Path("data/exports")
        self.export_dir.mkdir(parents=True, exist_ok=True)
        
        self.tasks: Dict[str, ExportTask] = {}
        self.max_tasks = 50

    def _prune_tasks(self):
        """限制内存任务与导出文件数量，防止匿名导出长期占用磁盘。"""
        if len(self.tasks) <= self.max_tasks:
            return
        expired = sorted(self.tasks.keys(), key=lambda tid: self.tasks[tid].create_time)[:-self.max_tasks]
        for task_id in expired:
            task = self.tasks.pop(task_id, None)
            if task and task.file_path:
                try:
                    Path(task.file_path).unlink(missing_ok=True)
                except OSError:
                    pass
    
    def _parse_markdown_content(self, content: str) -> List[Dict[str, Any]]:
        """解析Markdown内容为章节结构"""
        chapters = []
        lines = content.split('\n')
        current_chapter = None
        current_content = []
        
        for line in lines:
            if line.startswith('## '):
                if current_chapter:
                    current_chapter['content'] = '\n'.join(current_content)
                    chapters.append(current_chapter)
                
                current_chapter = {
                    'title': line[3:].strip(),
                    'content': '',
                    'sections': []
                }
                current_content = []
            elif line.startswith('### '):
                if current_content and current_chapter:
                    current_chapter['content'] = '\n'.join(current_content)
                    current_content = []
                
                if current_chapter:
                    current_chapter['sections'].append({
                        'title': line[4:].strip(),
                        'content': ''
                    })
            else:
                if current_chapter and current_chapter.get('sections'):
                    current_chapter['sections'][-1]['content'] += line + '\n'
                else:
                    current_content.append(line)
        
        if current_chapter:
            current_chapter['content'] = '\n'.join(current_content)
            chapters.append(current_chapter)
        
        return chapters
    
    def _build_report_data(self, report_type: ReportType,
                            chapters: List[Dict[str, Any]],
                            metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        """由章节结构构建报告数据（solution / competitor 共用）"""
        metadata = metadata or {}
        if report_type == ReportType.SOLUTION:
            return {
                'title': metadata.get('title', '华为云解决方案建议书'),
                'subtitle': '智能匹配生成报告',
                'create_date': datetime.now().strftime('%Y年%m月%d日'),
                'customer_name': metadata.get('customer', ''),
                'chapters': chapters,
                'appendix': []
            }
        competitor = metadata.get('competitor', '竞争对手')
        industry = metadata.get('industry', '行业')
        return {
            'title': f'华为云 vs {competitor} 竞争分析报告',
            'subtitle': f'{industry}行业',
            'create_date': datetime.now().strftime('%Y年%m月%d日'),
            'customer_name': '',
            'chapters': chapters,
            'appendix': []
        }

    def _render_and_save(self, report_data: Dict[str, Any],
                         report_type: ReportType,
                         format: ExportFormat,
                         file_prefix: str,
                         cost_reference: Dict[str, Any] = None) -> ExportTask:
        """根据报告数据渲染并保存为 Word/PDF，返回完成任务"""
        self._prune_tasks()
        task = ExportTask(format=format, report_type=report_type)
        self.tasks[task.task_id] = task
        try:
            task.status = TaskStatus.PROCESSING
            start_time = time.time()
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

            if format == ExportFormat.WORD:
                file_name = f"{file_prefix}_{timestamp}.docx"
                file_path = self.export_dir / file_name
                generator = WordGenerator()
                generator.generate_report(report_data)
                generator.save(str(file_path))
            elif format == ExportFormat.PPTX:
                file_name = f"{file_prefix}_{timestamp}.pptx"
                file_path = self.export_dir / file_name
                self._generate_pptx(report_data, str(file_path))
            else:
                file_name = f"{file_prefix}_{timestamp}.pdf"
                file_path = self.export_dir / file_name
                self._generate_pdf(report_data, str(file_path), cost_reference=cost_reference)

            file_size = file_path.stat().st_size

            task.status = TaskStatus.COMPLETED
            task.complete_time = datetime.now()
            task.file_path = str(file_path)
            task.file_name = file_name
            task.file_size = file_size
            task.download_url = f"/api/export/download/{task.task_id}"
        except Exception as e:
            task.status = TaskStatus.FAILED
            task.error_message = str(e)
        return task

    def generate_report(self, report_type: ReportType, 
                        content: str, 
                        format: ExportFormat = ExportFormat.WORD,
                        metadata: Dict[str, Any] = None,
                        cost_reference: Dict[str, Any] = None) -> ExportTask:
        """生成报告（输入为 Markdown 内容，内部解析为章节）"""
        metadata = metadata or {}
        if report_type == ReportType.SOLUTION:
            chapters = self._parse_markdown_content(content)
            if not chapters:
                chapters = [{'title': '解决方案分析', 'content': content, 'sections': []}]
            file_prefix = "solution_report"
        else:
            chapters = self._parse_markdown_content(content)
            if not chapters:
                chapters = [{'title': '竞争分析', 'content': content, 'sections': []}]
            file_prefix = "competitor_report"

        report_data = self._build_report_data(report_type, chapters, metadata)
        self._attach_cost_reference(report_data, cost_reference)
        return self._render_and_save(report_data, report_type, format, file_prefix, cost_reference=cost_reference)

    def generate_report_from_json(self, report_type: ReportType,
                                  chapters: List[Dict[str, Any]],
                                  format: ExportFormat = ExportFormat.WORD,
                                  metadata: Dict[str, Any] = None,
                                  cost_reference: Dict[str, Any] = None) -> ExportTask:
        """
        直接由结构化 JSON（章节+要点）生成报告。
        用于『一键出可编辑售前方案书』：方案生成器已产出 solution_json，
        此处直接消费，跳过 Markdown 重新解析，结构更精准。
        """
        metadata = metadata or {}
        if not chapters:
            chapters = [{'title': '解决方案分析', 'content': '', 'sections': []}]
        file_prefix = "solution_report" if report_type == ReportType.SOLUTION else "competitor_report"
        report_data = self._build_report_data(report_type, chapters, metadata)
        self._attach_cost_reference(report_data, cost_reference)
        return self._render_and_save(report_data, report_type, format, file_prefix, cost_reference=cost_reference)

    # ---------------------------------------------------------------
    # 成本参考附表（导出时由前端成本卡片编辑态传入）
    # 关键：以与原方案内『对比表格』完全一致的「原生表格渲染」输出，
    # 严禁将价格表拼成 Markdown 字符串塞进正文（会造成 Word/PDF 乱码）。
    # Word 走 _render_markdown -> _add_table（python-docx Table Grid）；
    # PDF 走 reportlab Table。两者均为原生表格。
    # ---------------------------------------------------------------
    def _attach_cost_reference(self, report_data: Dict[str, Any], cost_reference: Dict[str, Any]):
        """把成本参考数据作为『附录章节』挂到 report_data，内容用 Markdown 表格描述，
        由 WordGenerator 以原生表格渲染（与方案内对比表格同一通道）。"""
        if not cost_reference:
            return
        cr_md = self._build_cost_reference_markdown(cost_reference)
        if not cr_md:
            return
        industry = cost_reference.get('industry', '') or ''
        title = f"成本参考估算（{industry or '通用'}行业 · 区间参考，非精确报价）"
        report_data.setdefault('appendix', [])
        report_data['appendix'].append({'title': title, 'content': cr_md})

    def _build_cost_reference_markdown(self, cost_reference: Dict[str, Any]) -> str:
        """把成本参考 rows 渲染为 Markdown 表格文本（供 Word 原生表格渲染）。

        支持 view_mode(month/year)：年视图下单价/小计/合计按 12 × annual_discount 折算，
        列头与单位切换为「元/年」；存储的 unit_price 仍为月度基准。
        """
        rows = cost_reference.get('rows') or []
        if not rows:
            return ''
        view_mode = (cost_reference.get('view_mode') or 'month')
        try:
            disc = float(cost_reference.get('annual_discount', 0.85) or 0.85)
        except (TypeError, ValueError):
            disc = 0.85
        is_year = (view_mode == 'year')
        factor = (12 * disc) if is_year else 1
        unit = '年' if is_year else '月'
        lines = []
        lines.append(f'| 产品 | 规格 | 计费方式 | 数量 | 单价(元/{unit}) | 小计(元/{unit}) |')
        lines.append('| :--- | :--- | :--- | :--- | :--- | :--- |')
        total = 0.0
        for r in rows:
            product = (r.get('product', '') or '')
            spec = (r.get('spec', '') or '')
            if r.get('business_only'):
                note = r.get('note') or '商务报价，请咨询华为云销售'
                lines.append(f'| **{product}** | {spec} | — | — | — | 商务定价：{note} |')
                continue
            if r.get('no_price'):
                note = r.get('note') or '参考价待补充'
                lines.append(f'| **{product}** | {spec} | — | — | — | 参考价待补充：{note} |')
                continue
            billing = (r.get('billing', '') or '')
            unit_label = (r.get('unit_label', '') or '')
            bill_txt = billing + (f'·{unit_label}' if unit_label else '')
            try:
                qty = float(r.get('qty') or 0)
                unit_price = float(r.get('unit_price') or 0)
            except (TypeError, ValueError):
                qty, unit_price = 0.0, 0.0
            subtotal = round(qty * unit_price * factor, 2)
            total += subtotal
            warn = '' if r.get('verified', True) else ' ⚠'
            lines.append(
                f'| **{product}**{warn} | {spec} | {bill_txt} | {_fmt_num(qty)} | '
                f'{_fmt_num(round(unit_price * factor, 2))} | {_fmt_num(subtotal)} |'
            )
        lines.append('')
        view_tag = f'年费(年付≈月×12×{disc})' if is_year else '月费'
        lines.append(f'**合计（估算，不含商务定价与待补充项）：¥{_fmt_num(round(total, 2))} 元/{unit}（{view_tag}）**')
        disclaimer = cost_reference.get('disclaimer', '')
        if disclaimer:
            lines.append('')
            lines.append(f'免责声明：{disclaimer}')
        return '\n'.join(lines)
    
    def _generate_pptx(self, report_data: Dict[str, Any], file_path: str):
        """生成 PPTX 演示稿（python-pptx）。

        版式：
        - 封面页：标题 / 副标题 / 客户 / 日期
        - 每章一页：章节标题 + 内容要点（Markdown 标记做轻量清洗）
        - 成本参考（附录）按要点列出
        说明：PPTX 不渲染 Markdown 表格（降级为要点列表），保证任何中文内容不丢失。
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError as e:
            raise Exception(f"PPTX 生成依赖缺失（python-pptx）: {e}")

        HUAWEI_RED = RGBColor(0xC7, 0x00, 0x0B)
        DARK = RGBColor(0x33, 0x33, 0x33)
        GRAY = RGBColor(0x66, 0x66, 0x66)

        def _clean_md(text: str) -> str:
            """轻量清理 Markdown 标记：标题符号/加粗/行首列表符/链接"""
            text = re.sub(r'^#{1,6}\s+', '', text)
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'\*(.+?)\*', r'\1', text)
            text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
            text = re.sub(r'^\s*[-*•·]\s+', '• ', text)
            text = text.replace('`', '')
            return text.strip()

        def _add_bullets(slide, items, size=16):
            """在正文占位符中追加要点（支持 '\n' 内嵌换行拆段）"""
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            first = True
            for it in items:
                for line in it.split('\n'):
                    line = _clean_md(line)
                    if not line.strip():
                        continue
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = line[:180]
                    p.font.size = Pt(size)
                    p.font.color.rgb = DARK
                    p.space_after = Pt(4)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ---- 封面页 ----
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = cover.shapes.title
        title_shape.text = report_data.get('title', '华为云解决方案')
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = HUAWEI_RED
        sub = cover.placeholders[1]
        sub.text = (report_data.get('subtitle', '智能匹配生成报告') or '')
        sub.text_frame.paragraphs[0].font.size = Pt(18)
        sub.text_frame.paragraphs[0].font.color.rgb = GRAY

        # 客户/日期补充行（封面正文追加）
        cover_body = cover.placeholders[1].text_frame
        extra = []
        if report_data.get('customer_name'):
            extra.append(f"客户：{report_data['customer_name']}")
        if report_data.get('create_date'):
            extra.append(f"日期：{report_data['create_date']}")
        for line in extra:
            p = cover_body.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY

        # ---- 章节页 ----
        for chapter in report_data.get('chapters', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = _clean_md(chapter.get('title', ''))[:60]
            # 子节（### 标题）作为段落前置标签
            items = []
            for sec in chapter.get('sections', []):
                sec_title = _clean_md(sec.get('title', ''))
                sec_content = sec.get('content', '')
                if sec_title:
                    items.append(f"【{sec_title}】")
                if sec_content:
                    items.append(sec_content)
            content = chapter.get('content', '')
            if content:
                items.append(content)
            if not items:
                items = ['（本节无要点内容）']
            _add_bullets(slide, items)

        # ---- 附录（成本参考等） ----
        for appendix in report_data.get('appendix', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = _clean_md(appendix.get('title', '附录'))[:60]
            _add_bullets(slide, [appendix.get('content', '')], size=14)

        prs.save(file_path)

    def _generate_pdf(self, report_data: Dict[str, Any], file_path: str,
                      cost_reference: Dict[str, Any] = None):
        """生成PDF报告（使用reportlab，成本参考以原生 Table 渲染，不乱码）"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib import colors
            from reportlab.lib.units import cm
            from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                           Table, TableStyle)
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            pdfmetrics.registerFont(TTFont('SimSun', 'simsun.ttc'))
            
            doc = SimpleDocTemplate(file_path, pagesize=A4)
            styles = getSampleStyleSheet()
            
            title_style = ParagraphStyle(
                'Title',
                parent=styles['Title'],
                fontName='SimSun',
                fontSize=24
            )
            
            body_style = ParagraphStyle(
                'Body',
                parent=styles['Normal'],
                fontName='SimSun',
                fontSize=12
            )
            
            cell_style = ParagraphStyle(
                'Cell',
                parent=styles['Normal'],
                fontName='SimSun',
                fontSize=9,
                leading=12
            )
            
            story = []
            
            story.append(Paragraph(report_data.get('title', ''), title_style))
            story.append(Spacer(1, 20))
            
            for chapter in report_data.get('chapters', []):
                story.append(Paragraph(chapter.get('title', ''), body_style))
                story.append(Spacer(1, 10))
                
                content = chapter.get('content', '')
                for para in content.split('\n'):
                    if not para.strip():
                        continue
                    # 列表项：去掉行首 [-*•·] 符号，与网页端无圆点风格一致（不编号）
                    cleaned = re.sub(r'^[-*•·]\s+', '', para.strip())
                    story.append(Paragraph(cleaned, body_style))
            
            # ===== 成本参考附表（原生表格，避免 Markdown 乱码） =====
            if cost_reference:
                cr_rows = cost_reference.get('rows') or []
                if cr_rows:
                    cr_view = (cost_reference.get('view_mode') or 'month')
                    try:
                        cr_disc = float(cost_reference.get('annual_discount', 0.85) or 0.85)
                    except (TypeError, ValueError):
                        cr_disc = 0.85
                    cr_is_year = (cr_view == 'year')
                    cr_factor = (12 * cr_disc) if cr_is_year else 1
                    cr_unit = '年' if cr_is_year else '月'
                    story.append(Spacer(1, 16))
                    industry = cost_reference.get('industry', '') or ''
                    story.append(Paragraph(
                        f'成本参考估算（{industry or "通用"}行业 · 区间参考，非精确报价）',
                        ParagraphStyle('CRTitle', parent=body_style, fontSize=13, spaceAfter=6)
                    ))
                    header = ['产品', '规格', '计费方式', '数量', f'单价(元/{cr_unit})', f'小计(元/{cr_unit})']
                    table_data = [[Paragraph(f'<b>{h}</b>', cell_style) for h in header]]
                    total = 0.0
                    for r in cr_rows:
                        product = (r.get('product', '') or '')
                        spec = (r.get('spec', '') or '')
                        if r.get('business_only'):
                            note = r.get('note') or '商务报价，请咨询华为云销售'
                            table_data.append([
                                Paragraph(f'<b>{product}</b>', cell_style), Paragraph(spec, cell_style),
                                Paragraph('—', cell_style), Paragraph('—', cell_style),
                                Paragraph('—', cell_style), Paragraph(f'商务定价：{note}', cell_style)
                            ])
                            continue
                        if r.get('no_price'):
                            note = r.get('note') or '参考价待补充'
                            table_data.append([
                                Paragraph(f'<b>{product}</b>', cell_style), Paragraph(spec, cell_style),
                                Paragraph('—', cell_style), Paragraph('—', cell_style),
                                Paragraph('—', cell_style), Paragraph(f'参考价待补充：{note}', cell_style)
                            ])
                            continue
                        billing = (r.get('billing', '') or '')
                        unit_label = (r.get('unit_label', '') or '')
                        bill_txt = billing + (f'·{unit_label}' if unit_label else '')
                        try:
                            qty = float(r.get('qty') or 0)
                            unit_price = float(r.get('unit_price') or 0)
                        except (TypeError, ValueError):
                            qty, unit_price = 0.0, 0.0
                        subtotal = round(qty * unit_price * cr_factor, 2)
                        total += subtotal
                        table_data.append([
                            Paragraph(product, cell_style), Paragraph(spec, cell_style),
                            Paragraph(bill_txt, cell_style), Paragraph(_fmt_num(qty), cell_style),
                            Paragraph(_fmt_num(round(unit_price * cr_factor, 2)), cell_style),
                            Paragraph(_fmt_num(subtotal), cell_style)
                        ])
                    cr_table = Table(table_data, repeatRows=1, colWidths=[
                        3.0 * cm, 4.2 * cm, 2.6 * cm, 1.4 * cm, 2.6 * cm, 2.8 * cm
                    ])
                    cr_table.setStyle(TableStyle([
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#999999')),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#C7000B')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F7F7F7')]),
                        ('LEFTPADDING', (0, 0), (-1, -1), 4),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                        ('TOPPADDING', (0, 0), (-1, -1), 3),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
                    ]))
                    story.append(cr_table)
                    story.append(Spacer(1, 6))
                    if cr_is_year:
                        view_tag = f'年费(年付≈月×12×{cr_disc})'
                    else:
                        view_tag = '月费'
                    story.append(Paragraph(
                        f'合计（估算，不含商务定价与待补充项）：¥{_fmt_num(round(total, 2))} 元/{cr_unit}（{view_tag}）',
                        ParagraphStyle('CRTotal', parent=body_style, fontSize=10, spaceAfter=4)
                    ))
                    disclaimer = cost_reference.get('disclaimer', '')
                    if disclaimer:
                        story.append(Paragraph(
                            f'免责声明：{disclaimer}',
                            ParagraphStyle('CRDisc', parent=body_style, fontSize=8, textColor=colors.HexColor('#666666'))
                        ))
            
            doc.build(story)
            
        except Exception as e:
            raise Exception(f"PDF生成失败: {str(e)}")
    
    def get_task(self, task_id: str) -> ExportTask:
        """获取任务状态"""
        return self.tasks.get(task_id)
    
    def get_file_path(self, task_id: str) -> str:
        """获取文件路径"""
        task = self.tasks.get(task_id)
        if task and task.status == TaskStatus.COMPLETED:
            return task.file_path
        return None
