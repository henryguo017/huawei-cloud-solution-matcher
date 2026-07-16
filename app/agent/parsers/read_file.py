"""
多格式文件解析 —— 把客户上传的任意格式文件提取为纯文本。

设计：
- 按扩展名分发到对应解析器（全部懒加载第三方库，缺失时友好报错，不影响其他格式）
- 图片走 OCR（见 ocr.py）
- PDF 文本为空（扫描件）自动回退到页面渲染 + OCR
- 分块全量提取：超长文本按重叠窗口切片，确保全部内容进入上下文、零丢弃（不盲截断）
"""
import os
import logging
from typing import List

logger = logging.getLogger(__name__)

# 支持的扩展名（与上传接口校验保持一致）
ALLOWED_EXT = {
    ".docx", ".xlsx", ".pdf", ".pptx",
    ".txt", ".csv", ".md",
    ".png", ".jpg", ".jpeg",
}

# 分块全量提取默认参数（用户拍板）
CHUNK_SIZE = 8000
CHUNK_OVERLAP = 500

# 扫描件 PDF 最多 OCR 前 N 页，避免过慢
PDF_OCR_MAX_PAGES = 20


# ============================================================
# 各格式解析器（懒加载依赖）
# ============================================================

def _read_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _read_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"### Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"### Slide {i + 1}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_pdf(path: str) -> str:
    try:
        import fitz  # pymupdf
    except ImportError:
        return "Error: 未安装 pymupdf，无法读取 PDF"
    try:
        doc = fitz.open(path)
        text = "\n".join(page.get_text() for page in doc)
        # 扫描件（文本为空）→ 回退 OCR
        if not text.strip():
            from app.agent.parsers.ocr import ocr_pil
            from PIL import Image
            import io
            ocr_parts = []
            for idx, page in enumerate(doc):
                if idx >= PDF_OCR_MAX_PAGES:
                    break
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes()))
                ocr_parts.append(ocr_pil(img))
            doc.close()
            return "\n".join(ocr_parts)
        doc.close()
        return text
    except Exception as e:
        return f"Error: 读取 PDF 失败: {e}"


def _read_image(path: str) -> str:
    from app.agent.parsers.ocr import ocr_image
    return ocr_image(path)


# ============================================================
# 统一入口
# ============================================================

def extract_text(path: str) -> str:
    """提取任意支持格式文件的纯文本。错误返回 Error: 前缀字符串（不抛异常）。"""
    ext = os.path.splitext(path)[1].lower()
    if ext not in ALLOWED_EXT:
        return f"Error: 不支持的文件格式 {ext}"
    try:
        if ext == ".docx":
            return _read_docx(path)
        elif ext == ".xlsx":
            return _read_xlsx(path)
        elif ext == ".pdf":
            return _read_pdf(path)
        elif ext == ".pptx":
            return _read_pptx(path)
        elif ext in (".txt", ".csv", ".md"):
            return _read_text(path)
        elif ext in (".png", ".jpg", ".jpeg"):
            return _read_image(path)
        return f"Error: 未处理的扩展名 {ext}"
    except Exception as e:
        return f"Error: 读取文件失败: {e}"


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """
    分块全量提取：把长文本按重叠窗口切片。

    返回单块表示无需切片；多块时保留全部内容（重叠避免句子被切断），零丢弃。
    """
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start = end - overlap
    return chunks
